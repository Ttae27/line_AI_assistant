import io
import fitz 
import openpyxl
from enum import Enum
from docx import Document
from pptx import Presentation
from tools import llm
from langchain_core.messages import HumanMessage

def extract_text_from_pdf_bytes(file_bytes: bytearray) -> str:
    file_stream = io.BytesIO(file_bytes)
    text = ""
    with fitz.open(stream=file_stream, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx_bytes(file_bytes: bytearray) -> str:
    file_stream = io.BytesIO(file_bytes)
    doc = Document(file_stream)
    full_text = []

    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                full_text.append(" | ".join(row_text))

    return "\n".join(full_text)

def extract_text_from_xlsx_bytes(file_bytes: bytearray) -> str:
    file_stream = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(file_stream, data_only=True)
    all_text = []

    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = [str(cell).strip() for cell in row if cell is not None]
            if row_text:
                all_text.append(" | ".join(row_text))

    return "\n".join(all_text)

def extract_text_from_pptx_bytes(file_bytes: bytearray) -> str:
    file_stream = io.BytesIO(file_bytes)
    prs = Presentation(file_stream)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content = shape.text.strip()
                if content:
                    text_runs.append(content)

    return "\n".join(text_runs)

def summarized(file_bytes: bytearray, file_type):
    if file_type == 'docx':
        text = extract_text_from_docx_bytes(file_bytes)
    elif file_type == 'xlsx':
        text = extract_text_from_xlsx_bytes(file_bytes)
    elif file_type == 'pdf':
        text = extract_text_from_pdf_bytes(file_bytes)
    elif file_type == 'pptx':
        text = extract_text_from_pptx_bytes(file_bytes)
    query = 'short summarization this document: ' + text
    result = llm.invoke([HumanMessage(query)])
    return result.content

